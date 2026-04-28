"""
doc_generator.py — Export LLM-generated analysis content to downloadable documents.

Supported formats
-----------------
  word   → .docx  (python-docx)
  excel  → .xlsx  (openpyxl)
  pptx   → .pptx  (python-pptx)
  pdf    → .pdf   (fpdf2)

Public API
----------
    generate_document(fmt, title, content) -> (bytes, filename, mime_type)
"""

from __future__ import annotations

import io
import re
from typing import Literal

DocFormat = Literal["word", "excel", "pptx", "pdf"]


# ---------------------------------------------------------------------------
# Markdown → section list parser
# ---------------------------------------------------------------------------

def _parse_sections(text: str) -> list[dict[str, str]]:
    """Parse markdown into a list of {heading, body} dicts."""
    sections: list[dict[str, str]] = []
    heading = "Overview"
    body: list[str] = []

    for line in text.split("\n"):
        if line.startswith(("### ", "## ", "# ")):
            if body:
                sections.append({"heading": heading, "body": "\n".join(body).strip()})
            # Strip leading #s
            heading = re.sub(r"^#+\s*", "", line).strip()
            body = []
        else:
            body.append(line)

    if body:
        sections.append({"heading": heading, "body": "\n".join(body).strip()})

    return [s for s in sections if s["heading"].strip() or s["body"].strip()]


def _clean_line(line: str) -> str:
    """Remove markdown bold/italic/code markers."""
    line = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
    line = re.sub(r"\*(.+?)\*",     r"\1", line)
    line = re.sub(r"`(.+?)`",       r"\1", line)
    return line.strip()


def _is_bullet(line: str) -> bool:
    return bool(re.match(r"^\s*[-*•]", line))


def _is_numbered(line: str) -> bool:
    return bool(re.match(r"^\s*\d+[.)]\s", line))


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------

def generate_word(title: str, content: str) -> bytes:
    """Return .docx bytes for *content* (markdown)."""
    from docx import Document                          # type: ignore[import]
    from docx.shared import Pt, RGBColor              # type: ignore[import]
    from docx.enum.text import WD_ALIGN_PARAGRAPH     # type: ignore[import]

    doc = Document()

    # Page title
    t = doc.add_heading(title, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER           # type: ignore[attr-defined]

    for section in _parse_sections(content):
        if section["heading"] and section["heading"] != "Overview":
            doc.add_heading(section["heading"], level=2)
        for line in section["body"].split("\n"):
            if not line.strip():
                continue
            clean = _clean_line(line)
            if _is_bullet(line):
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(re.sub(r"^\s*[-*•]\s*", "", clean))
            elif _is_numbered(line):
                p = doc.add_paragraph(style="List Number")
                p.add_run(re.sub(r"^\s*\d+[.)]\s*", "", clean))
            else:
                doc.add_paragraph(clean)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def generate_excel(title: str, content: str) -> bytes:
    """Return .xlsx bytes with content organised in a structured sheet."""
    from openpyxl import Workbook                      # type: ignore[import]
    from openpyxl.styles import Font, PatternFill, Alignment  # type: ignore[import]

    wb = Workbook()
    ws = wb.active
    ws.title = "Analysis"

    # ── Title row ──
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor="1F4E79")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws.merge_cells("A1:C1")
    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 65
    ws.column_dimensions["C"].width = 20

    row = 3
    for section in _parse_sections(content):
        # Section heading row
        cell = ws.cell(row=row, column=1, value=section["heading"])
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="2E75B6")
        ws.merge_cells(f"A{row}:C{row}")
        row += 1
        # Body lines
        for line in section["body"].split("\n"):
            if not line.strip():
                continue
            clean = _clean_line(re.sub(r"^\s*[-*•\d.)]+\s*", "", line))
            c = ws.cell(row=row, column=2, value=clean)
            c.alignment = Alignment(wrap_text=True)
            row += 1
        row += 1  # blank separator between sections

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def generate_pptx(title: str, content: str) -> bytes:
    """Return .pptx bytes with one slide per section."""
    from pptx import Presentation                     # type: ignore[import]
    from pptx.util import Inches, Pt                  # type: ignore[import]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Generated by advanced-QODE | QODE Framework Analysis"

    for section in _parse_sections(content):
        body = section["body"].strip()
        if not body:
            continue
        slide  = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section["heading"]
        tf     = slide.placeholders[1].text_frame
        tf.clear()
        first  = True
        for line in body.split("\n"):
            if not line.strip():
                continue
            clean = _clean_line(re.sub(r"^\s*[-*•\d.)]+\s*", "", line))
            p     = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text  = clean
            p.level = 1 if (_is_bullet(line) or _is_numbered(line)) else 0
            first   = False

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def generate_pdf(title: str, content: str) -> bytes:
    """Return PDF bytes using fpdf2."""
    from fpdf import FPDF  # type: ignore[import]

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Title banner
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_fill_color(31, 78, 121)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 12, _safe_ascii(title), fill=True, new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(4)

    for section in _parse_sections(content):
        if section["heading"] and section["heading"] != "Overview":
            pdf.set_font("Helvetica", "B", 12)
            pdf.set_fill_color(46, 117, 182)
            pdf.set_text_color(255, 255, 255)
            pdf.cell(0, 8, _safe_ascii(section["heading"]),
                     fill=True, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(30, 30, 30)
        for line in section["body"].split("\n"):
            if not line.strip():
                continue
            clean = _safe_ascii(_clean_line(line))
            if _is_bullet(line):
                pdf.cell(6, 6, chr(149), new_x="RIGHT", new_y="TOP")
                pdf.multi_cell(0, 6, re.sub(r"^\s*[-*•]\s*", "", clean))
            elif _is_numbered(line):
                pdf.multi_cell(0, 6, clean)
            else:
                pdf.multi_cell(0, 6, clean)
        pdf.ln(3)

    return bytes(pdf.output())


def _safe_ascii(text: str) -> str:
    """Replace non-latin-1 characters so fpdf2 (latin-1 font) won't crash."""
    try:
        text.encode("latin-1")
        return text
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text.encode("ascii", "replace").decode("ascii")


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------

_MIME: dict[str, str] = {
    "word":  "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx":  "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf":   "application/pdf",
}
_EXT: dict[str, str] = {"word": "docx", "excel": "xlsx", "pptx": "pptx", "pdf": "pdf"}


def generate_document(
    fmt: DocFormat,
    title: str,
    content: str,
) -> tuple[bytes, str, str]:
    """Generate a document in the requested format.

    Args:
        fmt:     One of "word", "excel", "pptx", "pdf".
        title:   Document title (used in heading and filename).
        content: Markdown-formatted text from the LLM.

    Returns:
        ``(file_bytes, filename, mime_type)``
    """
    safe = re.sub(r"[^\w\s-]", "", title)[:50].strip().replace(" ", "_") or "QODE_Report"
    filename = f"{safe}.{_EXT[fmt]}"

    generators = {
        "word":  generate_word,
        "excel": generate_excel,
        "pptx":  generate_pptx,
        "pdf":   generate_pdf,
    }
    if fmt not in generators:
        raise ValueError(f"Unsupported format: {fmt!r}")

    data = generators[fmt](title, content)
    return data, filename, _MIME[fmt]
